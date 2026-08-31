import io
import re
from pathlib import Path

import fitz
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def clean_text(value: str) -> str:
    value = (value or "").replace("\x00", " ").replace("\u00ad", "")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def is_noise(text: str) -> bool:
    text = clean_text(text)
    if len(text) < 45:
        return True
    lower = text.lower()
    urls = len(re.findall(r"https?://|www\.", lower))
    citations = len(re.findall(r"\[\d+\]|\b(?:isbn|doi|retrieved|accessed|archived)\b", lower))
    digits = len(re.findall(r"\d", text))
    return urls > 0 or citations >= 2 or digits > len(text) * 0.24


def section(location: str, text: str) -> dict | None:
    text = clean_text(text)
    return {"location": location, "text": text} if text and not is_noise(text) else None


def extract_pdf(data: bytes) -> tuple[list[dict], int]:
    pdf = fitz.open(stream=data, filetype="pdf")
    sections, image_count = [], 0
    try:
        for page_no, page in enumerate(pdf, start=1):
            image_count += len(page.get_images(full=True))
            page_sections = []
            for block_no, block in enumerate(page.get_text("blocks", sort=True), start=1):
                item = section(f"Sayfa {page_no}, blok {block_no}", block[4] if len(block) > 4 else "")
                if item:
                    page_sections.append(item)
            if page_sections:
                sections.extend(page_sections)
            else:
                item = section(f"Sayfa {page_no}", page.get_text("text", sort=True))
                if item:
                    sections.append(item)
    finally:
        pdf.close()
    return sections, image_count


def extract_text(data: bytes) -> list[dict]:
    for encoding in ("utf-8-sig", "utf-8", "cp1254", "latin-1"):
        try:
            item = section("Dosya", data.decode(encoding))
            return [item] if item else []
        except UnicodeDecodeError:
            pass
    raise ValueError("Metin dosyasının karakter kodlaması okunamadı.")


def extract_docx(data: bytes) -> list[dict]:
    doc = Document(io.BytesIO(data))
    parts = [clean_text(p.text) for p in doc.paragraphs if clean_text(p.text)]
    for table_no, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            values = [clean_text(cell.text) for cell in row.cells]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            parts.append(f"Tablo {table_no}:\n" + "\n".join(rows))
    item = section("Word dokümanı", "\n\n".join(parts))
    return [item] if item else []


def extract_csv(data: bytes) -> list[dict]:
    last_error = None
    for encoding in ("utf-8-sig", "utf-8", "cp1254", "latin-1"):
        for separator in (None, ",", ";", "\t"):
            try:
                frame = pd.read_csv(io.BytesIO(data), encoding=encoding, sep=separator, engine="python").fillna("")
                item = section("CSV tablo", frame.to_csv(index=False))
                return [item] if item else []
            except Exception as error:
                last_error = error
    raise ValueError(f"CSV okunamadı: {last_error}")


def extract_xlsx(data: bytes) -> list[dict]:
    workbook = load_workbook(io.BytesIO(data), data_only=True)
    result = []
    for worksheet in workbook.worksheets:
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                rows.append(" | ".join(values))
        item = section(f"Excel sayfası: {worksheet.title}", "\n".join(rows))
        if item:
            result.append(item)
    return result


def extract_pptx(data: bytes) -> list[dict]:
    presentation = Presentation(io.BytesIO(data))
    result = []
    for slide_no, slide in enumerate(presentation.slides, start=1):
        parts = [clean_text(shape.text) for shape in slide.shapes if hasattr(shape, "text") and clean_text(shape.text)]
        item = section(f"Slayt {slide_no}", "\n".join(parts))
        if item:
            result.append(item)
    return result


def extract_html(data: bytes) -> list[dict]:
    soup = BeautifulSoup(data.decode("utf-8", errors="replace"), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    item = section("HTML dokümanı", soup.get_text("\n"))
    return [item] if item else []


def extract_document(filename: str, data: bytes) -> tuple[list[dict], dict]:
    extension = Path(filename).suffix.lower().lstrip(".")
    metadata = {"name": filename, "extension": extension, "size_bytes": len(data), "image_count": 0, "ocr_required": False}
    if extension == "pdf":
        sections, images = extract_pdf(data)
        metadata["image_count"] = images
        metadata["ocr_required"] = not bool(sections)
    elif extension in {"txt", "md"}:
        sections = extract_text(data)
    elif extension == "docx":
        sections = extract_docx(data)
    elif extension == "csv":
        sections = extract_csv(data)
    elif extension == "xlsx":
        sections = extract_xlsx(data)
    elif extension == "pptx":
        sections = extract_pptx(data)
    elif extension in {"html", "htm"}:
        sections = extract_html(data)
    else:
        raise ValueError(f"Desteklenmeyen dosya uzantısı: .{extension}")
    return sections, metadata
